"""Fallback slide viewer when PowerPoint is not installed.

Uses python-pptx to extract slide content (text, images) and displays
them in a simple fullscreen PySide6 window. No animations are supported
---this is a degraded experience for when PowerPoint is unavailable.
"""

from pathlib import Path
from io import BytesIO
import logging

from PySide6.QtCore import Qt, QSize, QRect, QPoint
from PySide6.QtWidgets import (
    QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QLabel,
    QPushButton, QGraphicsView, QGraphicsScene, QGraphicsPixmapItem,
    QScrollArea, QFrame,
)
from PySide6.QtGui import (
    QPainter, QColor, QFont, QPixmap, QImage, QPen, QBrush, QFontMetrics,
    QMouseEvent, QKeyEvent,
)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from PIL import Image

logger = logging.getLogger(__name__)


def _emu_to_px(emu: int, dpi: int = 96) -> float:
    """Convert EMU (English Metric Units) to pixels at given DPI."""
    # 1 EMU = 1/914400 inch
    return emu * dpi / 914400


class SlideWidget(QFrame):
    """Renders a single slide's content."""

    def __init__(self, slide, slide_width_px: int, slide_height_px: int, parent=None):
        super().__init__(parent)
        self.setFixedSize(slide_width_px, slide_height_px)
        self.setStyleSheet("background: white; border: 1px solid #ccc;")
        self._slide = slide
        self._width_px = slide_width_px
        self._height_px = slide_height_px

    def paintEvent(self, event) -> None:
        """Render slide content."""
        super().paintEvent(event)
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)

        slide_width_emu = self._slide.slide_width  # in EMU
        slide_height_emu = self._slide.slide_height

        scale_x = self._width_px / slide_width_emu if slide_width_emu else 1
        scale_y = self._height_px / slide_height_emu if slide_height_emu else 1

        for shape in self._slide.shapes:
            self._render_shape(painter, shape, scale_x, scale_y)

    def _render_shape(self, painter, shape, scale_x, scale_y) -> None:
        """Render a single shape."""
        left_px = int(shape.left * scale_x) if shape.left else 0
        top_px = int(shape.top * scale_y) if shape.top else 0
        width_px = int(shape.width * scale_x) if shape.width else 100
        height_px = int(shape.height * scale_y) if shape.height else 50

        # Check for text
        if shape.has_text_frame:
            self._render_text(painter, shape, left_px, top_px, width_px, height_px)

        # Check for image
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            self._render_image(painter, shape, left_px, top_px, width_px, height_px)

    def _render_text(self, painter, shape, x, y, w, h) -> None:
        """Render text from a shape."""
        for para in shape.text_frame.paragraphs:
            text = para.text
            if not text.strip():
                continue

            # Determine font
            font_size = 12
            font_name = "Segoe UI"
            is_bold = False
            color = QColor(0, 0, 0)

            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size.pt
                if run.font.name:
                    font_name = run.font.name
                if run.font.bold:
                    is_bold = True
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    color = QColor(
                        (rgb >> 16) & 0xFF,
                        (rgb >> 8) & 0xFF,
                        rgb & 0xFF,
                    )

            font = QFont(font_name, int(font_size))
            font.setBold(is_bold)
            painter.setFont(font)
            painter.setPen(color)

            align = para.alignment
            flags = Qt.AlignLeft | Qt.AlignTop | Qt.TextWordWrap
            if align == PP_ALIGN.CENTER:
                flags = Qt.AlignHCenter | Qt.AlignTop | Qt.TextWordWrap
            elif align == PP_ALIGN.RIGHT:
                flags = Qt.AlignRight | Qt.AlignTop | Qt.TextWordWrap

            rect = QRect(x, y, w, h)
            painter.drawText(rect, flags, text)
            y += int(font_size * 1.5)  # Move down for next paragraph

    def _render_image(self, painter, shape, x, y, w, h) -> None:
        """Render an embedded image."""
        try:
            image = shape.image
            if image:
                img_bytes = image.blob
                img = Image.open(BytesIO(img_bytes))
                img = img.convert("RGBA")
                data = img.tobytes("raw", "RGBA")
                qimage = QImage(
                    data, img.width, img.height,
                    img.width * 4, QImage.Format_RGBA8888,
                )
                pixmap = QPixmap.fromImage(qimage)
                scaled = pixmap.scaled(
                    w, h, Qt.KeepAspectRatio, Qt.SmoothTransformation
                )
                painter.drawPixmap(x, y, scaled)
        except Exception as e:
            logger.debug(f"Failed to render image: {e}")


class FallbackViewer(QMainWindow):
    """Simple fullscreen slide viewer using python-pptx.

    No animations. Static content only (text + images).
    Touch-friendly navigation buttons.
    """

    def __init__(self, file_path: str):
        super().__init__()
        self._file_path = file_path
        self._prs = Presentation(file_path)
        self._current_idx = 0
        self._total_slides = len(self._prs.slides)

        self._setup_ui()
        self._render_slide(0)

    def _setup_ui(self) -> None:
        """Set up fullscreen window with navigation."""
        self.setWindowTitle("PPT 简易查看器 (无动画)")
        self.setWindowState(Qt.WindowFullScreen)
        self.setStyleSheet("background: #1a1a1a;")

        central = QWidget(self)
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)
        layout.setContentsMargins(0, 0, 0, 10)

        # Slide info
        self._info_label = QLabel(f"1 / {self._total_slides}")
        self._info_label.setAlignment(Qt.AlignCenter)
        self._info_label.setStyleSheet(
            "color: #aaa; font-size: 14px; padding: 4px;"
        )
        layout.addWidget(self._info_label)

        # Scroll area for slide content
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setStyleSheet("background: #222; border: none;")
        scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)

        self._slide_container = QWidget()
        self._slide_container.setStyleSheet("background: transparent;")
        self._slide_layout = QVBoxLayout(self._slide_container)
        self._slide_layout.setAlignment(Qt.AlignCenter)

        scroll.setWidget(self._slide_container)
        layout.addWidget(scroll, 1)

        # Navigation buttons
        btn_layout = QHBoxLayout()
        btn_layout.setSpacing(20)

        self._prev_btn = QPushButton("◀ 上一页")
        self._prev_btn.setFixedSize(120, 60)
        self._prev_btn.clicked.connect(self._prev_slide)

        self._next_btn = QPushButton("下一页 ▶")
        self._next_btn.setFixedSize(120, 60)
        self._next_btn.clicked.connect(self._next_slide)

        self._exit_btn = QPushButton("✕ 退出")
        self._exit_btn.setFixedSize(100, 60)
        self._exit_btn.clicked.connect(self.close)

        btn_style = """
            QPushButton {
                background: rgba(0, 0, 0, 180);
                color: white;
                border: 2px solid #555;
                border-radius: 12px;
                font-size: 18px;
                font-weight: bold;
            }
            QPushButton:hover { background: rgba(0, 120, 212, 200); }
            QPushButton:pressed { background: rgba(0, 80, 160, 220); }
        """
        self._prev_btn.setStyleSheet(btn_style)
        self._next_btn.setStyleSheet(btn_style)
        self._exit_btn.setStyleSheet(btn_style)

        btn_layout.addStretch()
        btn_layout.addWidget(self._prev_btn)
        btn_layout.addWidget(self._next_btn)
        btn_layout.addWidget(self._exit_btn)
        btn_layout.addStretch()
        layout.addLayout(btn_layout)

    def _render_slide(self, index: int) -> None:
        """Render a slide by index."""
        # Clear existing
        while self._slide_layout.count():
            item = self._slide_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

        if 0 <= index < self._total_slides:
            slide = self._prs.slides[index]

            # Calculate slide dimensions
            sw = slide.slide_width or 9144000  # 10 inches default
            sh = slide.slide_height or 5143500  # 7.5 inches default

            # Target display size
            target_w = 960
            target_h = int(target_w * sh / sw) if sw else 540

            slide_widget = SlideWidget(slide, target_w, target_h)
            self._slide_layout.addWidget(slide_widget)
            self._info_label.setText(f"{index + 1} / {self._total_slides}")

            # Update button states
            self._prev_btn.setEnabled(index > 0)
            self._next_btn.setEnabled(index < self._total_slides - 1)

    def _next_slide(self) -> None:
        """Go to next slide."""
        if self._current_idx < self._total_slides - 1:
            self._current_idx += 1
            self._render_slide(self._current_idx)

    def _prev_slide(self) -> None:
        """Go to previous slide."""
        if self._current_idx > 0:
            self._current_idx -= 1
            self._render_slide(self._current_idx)

    def keyPressEvent(self, event: QKeyEvent) -> None:
        """Keyboard navigation."""
        if event.key() == Qt.Key_Right or event.key() == Qt.Key_Space:
            self._next_slide()
        elif event.key() == Qt.Key_Left or event.key() == Qt.Key_Backspace:
            self._prev_slide()
        elif event.key() == Qt.Key_Escape:
            self.close()
        else:
            super().keyPressEvent(event)
